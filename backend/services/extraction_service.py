import io
from typing import List, Tuple

# pypdf is pure-Python, safe to import at startup (used for native text extraction, no Tesseract needed)
from pypdf import PdfReader


def _clean_text(value: str) -> str:
    return "\n".join(line.rstrip() for line in value.splitlines()).strip()


def _preprocess_for_ocr(image):
    # Lazy import — only executed when OCR is actually triggered
    from PIL import ImageOps
    gray = ImageOps.grayscale(image)
    return ImageOps.autocontrast(gray)


def extract_text_from_image_bytes(image_bytes: bytes, lang: str = "eng") -> str:
    try:
        from PIL import Image
        import pytesseract
    except ImportError as e:
        raise RuntimeError(f"OCR library not installed: {e}") from e

    image = Image.open(io.BytesIO(image_bytes))
    image = _preprocess_for_ocr(image)
    text = pytesseract.image_to_string(image, lang=lang, config="--psm 6")
    return _clean_text(text)


def extract_text_from_pdf_native(pdf_bytes: bytes, max_pages: int = 50) -> Tuple[str, int]:
    reader = PdfReader(io.BytesIO(pdf_bytes))
    page_count = min(len(reader.pages), max_pages)
    chunks: List[str] = []
    for i in range(page_count):
        chunks.append(reader.pages[i].extract_text() or "")
    return _clean_text("\n\n".join(chunks)), page_count


def extract_text_from_pdf_ocr(
    pdf_bytes: bytes, max_pages: int = 20, lang: str = "eng"
) -> Tuple[str, int]:
    try:
        import pypdfium2 as pdfium
        import pytesseract
        from PIL import Image  # noqa: F401 — needed by pdfium .to_pil()
    except ImportError as e:
        raise RuntimeError(f"OCR library not installed: {e}") from e

    pdf = pdfium.PdfDocument(pdf_bytes)
    page_count = min(len(pdf), max_pages)
    chunks: List[str] = []
    for i in range(page_count):
        page = pdf[i]
        pil_image = page.render(scale=2).to_pil()
        processed = _preprocess_for_ocr(pil_image)
        chunks.append(pytesseract.image_to_string(processed, lang=lang, config="--psm 6"))
        page.close()
    return _clean_text("\n\n".join(chunks)), page_count


def extract_text_from_docx(file_bytes: bytes) -> str:
    """Extract plain text from a DOCX file using mammoth."""
    import mammoth
    result = mammoth.extract_raw_text(io.BytesIO(file_bytes))
    return _clean_text(result.value)


def extract_text_from_pptx(file_bytes: bytes) -> str:
    """Extract plain text from a PPTX file using python-pptx."""
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    chunks = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs).strip()
                    if line:
                        chunks.append(line)
    return _clean_text("\n".join(chunks))


def extract_text_from_file(file_bytes: bytes, filename: str, content_type: str) -> str:
    """Extract raw text from a PDF, DOCX, PPTX, plain-text, or image file."""
    lower = filename.lower()
    if content_type == "text/plain" or lower.endswith(".txt"):
        return _clean_text(file_bytes.decode("utf-8", errors="replace"))
    if content_type == "application/pdf" or lower.endswith(".pdf"):
        try:
            text, _ = extract_text_from_pdf_native(file_bytes)
            if len(text) >= 50:
                return text
        except Exception:
            pass
        text, _ = extract_text_from_pdf_ocr(file_bytes)
        return text
    if lower.endswith(".docx") or content_type in (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ):
        return extract_text_from_docx(file_bytes)
    if lower.endswith(".pptx") or content_type in (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ):
        return extract_text_from_pptx(file_bytes)
    else:
        return extract_text_from_image_bytes(file_bytes)
